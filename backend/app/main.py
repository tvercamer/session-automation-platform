from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from pathlib import Path
from typing import List, Dict, Any, Set, Optional
import requests
import uvicorn
import json
import sys
import os
import re
import shutil
from pptx import Presentation

# --- IMPORTS FROM LIBRARY.PY ---
from library import (
    scan_directory,
    resolve_dropped_item,
    list_translatable_folders,
    get_translations,
    save_translations
)

# --- CONFIGURATION ---
if sys.platform == "win32":
    APP_DATA_ROOT = Path(os.environ["APPDATA"])
    SETTINGS_DIR = APP_DATA_ROOT / "Datawijs-SAP"
else:
    SETTINGS_DIR = Path.home() / ".datawijs-sap"

SETTINGS_FILE = SETTINGS_DIR / "settings.json"
SETTINGS_DIR.mkdir(parents=True, exist_ok=True)

# --- APP SETUP ---
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- DATA MODELS ---

class KeyLabel(BaseModel):
    code: str
    label: str
    matches: List[str] = []

class SettingsModel(BaseModel):
    library_path: str
    output_path: str
    hubspot_api_key: str = ""
    languages: List[KeyLabel] = []
    industries: List[KeyLabel] = []

class ResolveRequest(BaseModel):
    path: str

class SectionRequest(BaseModel):
    title: str
    topics: List[str]

class GenerateRequest(BaseModel):
    session_name: str
    date: str
    customer_name: str
    customer_industry: str
    industry_code: str
    language_code: str
    sections: List[SectionRequest]

class TransListPayload(BaseModel):
    rootPath: str

class TransLoadPayload(BaseModel):
    targetPath: str

class TransSavePayload(BaseModel):
    targetPath: str
    entries: Dict[str, Any]

# --- HELPERS ---

def get_ignored_terms() -> Set[str]:
    if not SETTINGS_FILE.exists(): return set()
    try:
        with open(SETTINGS_FILE, "r") as f:
            data = json.load(f)
            langs = {l['code'].lower() for l in data.get('languages', [])}
            inds = {i['code'].lower() for i in data.get('industries', [])}
            combined = langs.union(inds)
            combined.update(['en', 'gen'])
            final_set = combined.union({t.upper() for t in combined})
            return final_set
    except:
        return set()

def find_best_matches(library_path: Path, topic: str, lang: str, ind: str) -> List[Path]:
    base_topic = Path(topic).stem
    found_files = []

    candidates = [
        f"{base_topic}_{lang}_{ind}",
        f"{base_topic}_{ind}_{lang}",
        f"{base_topic}_{lang}",
        f"{base_topic}_{ind}",
        base_topic
    ]

    best_file = None
    for root, _, files in os.walk(library_path):
        for cand in candidates:
            for f in files:
                f_path = Path(root) / f
                if f_path.stem == cand:
                    best_file = f_path
                    break
            if best_file: break
        if best_file: break

    if best_file:
        found_files.append(best_file)
        solution_name = f"{best_file.stem}_solution{best_file.suffix}"
        solution_path = best_file.parent / solution_name
        if solution_path.exists():
            found_files.append(solution_path)

    return found_files

# --- TRANSLATION LOGIC (AANGEPAST AAN JOUW STRUCTUUR) ---

def load_json_safely(path: Path) -> Dict:
    if not path.exists():
        return {}
    try:
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # print(f"DEBUG: JSON loaded from {path.name}")
            return data
    except Exception as e:
        print(f"ERROR: Corrupte JSON in {path}: {e}")
        return {}

def flatten_translation(data: Dict, lang: str, industry: str) -> Dict[str, str]:
    """
    Verwerkt JSON structuur: Key -> Default/Industries -> Taal
    """
    result = {}

    # Normaliseer inputs (bv 'nl' -> 'NL')
    target_lang = lang.upper().strip()
    target_ind = industry.lower().strip()

    # Loop door elke translation key (bv "HelloWorld", "welcome")
    for key, content in data.items():
        if not isinstance(content, dict):
            continue # Skip foute structuren

        value_found = None

        # STAP 1: Check Specifieke Industrie
        if 'industries' in content and isinstance(content['industries'], dict):
            # Zoek case-insensitive naar industry key
            ind_block = None
            for k, v in content['industries'].items():
                if k.lower() == target_ind:
                    ind_block = v
                    break

            if ind_block:
                # Probeer de taal in deze industry block
                # We checken case-insensitive op keys (NL, nl, En, EN)
                for l_key, l_val in ind_block.items():
                    if l_key.upper() == target_lang:
                        value_found = l_val
                        break

                # Fallback naar EN binnen industry als NL niet bestaat
                if not value_found:
                    for l_key, l_val in ind_block.items():
                        if l_key.upper() == 'EN':
                            value_found = l_val
                            break

        # STAP 2: Check Default (als nog niks gevonden in industry)
        if not value_found and 'default' in content:
            def_block = content['default']

            # Probeer target taal
            for l_key, l_val in def_block.items():
                if l_key.upper() == target_lang:
                    value_found = l_val
                    break

            # Fallback naar EN
            if not value_found:
                for l_key, l_val in def_block.items():
                    if l_key.upper() == 'EN':
                        value_found = l_val
                        break

        # STAP 3: Opslaan
        if value_found:
            result[key] = value_found
            # print(f"DEBUG: Key '{key}' vertaald naar: '{value_found}'")
        else:
            # print(f"DEBUG: Geen vertaling gevonden voor '{key}' (Lang: {target_lang}, Ind: {target_ind})")
            pass

    return result

def get_file_context(library_root: Path, file_path: Path, lang: str, ind: str, base_vars: Dict) -> Dict[str, str]:
    """
    Bouwt de dictionary op: Root JSON -> Local JSON -> System Vars.
    """
    context = {}

    # 1. Global Translations (Library Root)
    global_trans_path = library_root / "translations.json"
    if global_trans_path.exists():
        raw_global = load_json_safely(global_trans_path)
        flat_global = flatten_translation(raw_global, lang, ind)
        context.update(flat_global)

    # 2. Local Translations (File folder)
    local_trans_path = file_path.parent / "translations.json"
    if local_trans_path.exists():
        raw_local = load_json_safely(local_trans_path)
        if raw_local:
            flat_local = flatten_translation(raw_local, lang, ind)
            context.update(flat_local)

    # 3. System Variables (Overrides alles)
    context.update(base_vars)

    return context

def apply_replacements(text_frame, context: Dict[str, str]):
    if not text_frame or not text_frame.text:
        return

    text = text_frame.text
    # Regex: [% key %] (met dash support)
    pattern = re.compile(r'\[%\s*([\w\-]+)\s*%\]')

    matches = pattern.findall(text)
    if not matches:
        return

    new_text = text
    replaced = False
    for key in matches:
        # Zoek case-insensitive in de context
        val = None
        if key in context:
            val = context[key]
        elif key.lower() in context:
            val = context[key.lower()]

        if val is not None:
            # Case-insensitive replace van de tag
            key_pattern = re.compile(r'\[%\s*' + re.escape(key) + r'\s*%\]', re.IGNORECASE)
            new_text = key_pattern.sub(str(val), new_text)
            replaced = True

    if replaced:
        try:
            text_frame.text = new_text
        except:
            pass

# --- ENDPOINTS ---

@app.get("/")
def read_root():
    return {"status": "online", "message": "SAP Backend"}

@app.get("/settings")
def get_settings():
    if not SETTINGS_FILE.exists():
        return {"library_path": "", "output_path": "", "languages": [], "industries": []}
    try:
        with open(SETTINGS_FILE, "r") as f:
            data = json.load(f)
            if "languages" not in data: data["languages"] = []
            if "industries" not in data: data["industries"] = []
            return data
    except Exception as e:
        return {"error": str(e)}

@app.post("/settings")
def save_settings(settings: SettingsModel):
    try:
        with open(SETTINGS_FILE, "w") as f:
            json.dump(settings.model_dump(), f, indent=4)
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/library")
def get_library():
    if not SETTINGS_FILE.exists(): return []
    try:
        with open(SETTINGS_FILE, "r") as f:
            data = json.load(f)
            path_str = data.get("library_path", "")
        if not path_str or not Path(path_str).exists(): return []

        ignored = get_ignored_terms()
        return scan_directory(Path(path_str), ignored)
    except:
        return []

@app.post("/library/resolve")
def resolve_drop(req: ResolveRequest):
    ignored = get_ignored_terms()
    return resolve_dropped_item(req.path, ignored)

@app.post("/session/generate")
async def generate_session(req: GenerateRequest):
    if not SETTINGS_FILE.exists():
        raise HTTPException(status_code=500, detail="Settings not found")

    try:
        with open(SETTINGS_FILE, "r") as f:
            settings = json.load(f)
            lib_path = Path(settings.get("library_path", ""))
            out_root = Path(settings.get("output_path", ""))

        # 1. Setup Folders
        folder_name = f"{req.date}_{req.customer_name}_{req.session_name}".replace(" ", "_")
        target_dir = out_root / folder_name
        target_dir.mkdir(parents=True, exist_ok=True)
        exercises_dir = target_dir / "exercises"
        exercises_dir.mkdir(exist_ok=True)

        print(f"INFO: Start generatie voor {req.customer_name}")

        # System vars
        system_vars = {
            "customer_name": req.customer_name,
            "customer": req.customer_name,
            "klant": req.customer_name,
            "session_name": req.session_name,
            "session": req.session_name,
            "date": req.date,
            "datum": req.date,
            "industry": req.customer_industry,
            "industry_code": req.industry_code,
            "language": req.language_code
        }

        files_copied = 0
        pptx_merge_list = []

        # --- VERZAMELEN ---
        # A. Intro
        intro_matches = find_best_matches(lib_path, "Intro", req.language_code, req.industry_code)
        if intro_matches: pptx_merge_list.append(intro_matches[0])

        # B. Sections
        for section in req.sections:
            for topic in section.topics:
                matches = find_best_matches(lib_path, topic, req.language_code, req.industry_code)
                for file_path in matches:
                    if file_path.suffix.lower() == '.pptx':
                        pptx_merge_list.append(file_path)
                    else:
                        dest_folder = exercises_dir
                        final_name = file_path.name
                        dest_path = dest_folder / final_name
                        if dest_path.exists():
                            safe_title = section.title.replace("/", "-").replace("\\", "-")
                            final_name = f"{safe_title} - {file_path.name}"
                            dest_path = dest_folder / final_name
                            counter = 1
                            while dest_path.exists():
                                final_name = f"{safe_title} - {file_path.stem}_{counter}{file_path.suffix}"
                                dest_path = dest_folder / final_name
                                counter += 1
                        try:
                            shutil.copy2(file_path, dest_path)
                            files_copied += 1
                        except: pass

        # C. Outro
        outro_matches = find_best_matches(lib_path, "Outro", req.language_code, req.industry_code)
        if outro_matches: pptx_merge_list.append(outro_matches[0])

        # --- MERGEN EN VERVANGEN ---
        if pptx_merge_list:
            print(f"INFO: Mergen van {len(pptx_merge_list)} files met placeholders...")

            master_path = pptx_merge_list[0]
            # Context voor Master laden
            master_context = get_file_context(lib_path, master_path, req.language_code, req.industry_code, system_vars)

            prs = Presentation(master_path)

            # 1. Master Slides updaten
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        apply_replacements(shape.text_frame, master_context)

            # 2. Append rest
            for i in range(1, len(pptx_merge_list)):
                sub_path = pptx_merge_list[i]

                # Context voor Sub file laden
                file_context = get_file_context(lib_path, sub_path, req.language_code, req.industry_code, system_vars)

                try:
                    sub_prs = Presentation(sub_path)
                    for slide in sub_prs.slides:
                        # Layout match
                        layout_idx = sub_prs.slide_layouts.index(slide.slide_layout)
                        try:
                            slide_layout = prs.slide_layouts[layout_idx]
                        except:
                            slide_layout = prs.slide_layouts[0]

                        new_slide = prs.slides.add_slide(slide_layout)

                        for shape in slide.shapes:
                            # 1. Placeholders
                            if shape.is_placeholder:
                                try:
                                    ph_idx = shape.placeholder_format.idx
                                    new_ph = new_slide.placeholders[ph_idx]
                                    if shape.has_text_frame:
                                        new_ph.text = shape.text_frame.text
                                        apply_replacements(new_ph.text_frame, file_context)
                                except KeyError: pass

                            # 2. Losse Shapes
                            else:
                                if shape.has_text_frame:
                                    new_shape = new_slide.shapes.add_textbox(
                                        shape.left, shape.top, shape.width, shape.height
                                    )
                                    new_shape.text = shape.text
                                    apply_replacements(new_shape.text_frame, file_context)

                except Exception as e:
                    print(f"ERROR processing {sub_path.name}: {e}")

            output_pptx = target_dir / "slides.pptx"
            prs.save(output_pptx)
            files_copied += 1

        return {
            "status": "success",
            "target_dir": str(target_dir),
            "files_count": files_copied
        }

    except Exception as e:
        print(f"FATAL: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# --- HUBSPOT & TRANSLATIONS (Keep as is) ---
@app.get("/hubspot/companies")
def get_hubspot_companies():
    if not SETTINGS_FILE.exists(): return []
    try:
        with open(SETTINGS_FILE, "r") as f:
            api_key = json.load(f).get("hubspot_api_key", "")
        if not api_key: return []

        url = "https://api.hubapi.com/crm/v3/objects/companies"
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        params = {"limit": 100, "properties": "name,industry", "sort": "name"}

        response = requests.get(url, headers=headers, params=params)
        response.raise_for_status()
        data = response.json()
        return [{"name": r["properties"].get("name"), "code": r["id"], "industry": r["properties"].get("industry", "")}
                for r in data.get("results", []) if r["properties"].get("name")]
    except: return []

@app.post("/library/translations/folders")
def get_trans_folders(req: TransListPayload):
    return list_translatable_folders(req.rootPath)

@app.post("/library/translations/load")
def load_trans(req: TransLoadPayload):
    return get_translations(req.targetPath)

@app.post("/library/translations/save")
def save_trans(req: TransSavePayload):
    if not save_translations(req.targetPath, req.entries):
        raise HTTPException(status_code=500, detail="Failed to save")
    return {"success": True}

if __name__ == "__main__":
    uvicorn.run(app, host="127.0.0.1", port=8000)