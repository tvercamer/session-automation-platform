import os
import re
import shutil
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation

# Import TranslationService to build context
try:
    from app.services.translation_service import TranslationService
except ImportError:
    try:
        from .translation_service import TranslationService
    except ImportError:
        import translation_service as TranslationService

class GeneratorService:
    """
    Handles the generation of the final session output:
    - Finds best matching files
    - Copies exercises
    - Merges PPTX files
    - Replaces placeholders (text & images)
    """

    @staticmethod
    def find_best_matches(library_path: Path, topic: str, lang: str, ind: str) -> List[Path]:
        """
        Finds the most specific file version (Topic_Lang_Ind) or falls back to generic.
        """
        base_topic = Path(topic).stem
        found_files = []

        # Priority: Specific > Lang > Ind > Generic
        candidates = [
            f"{base_topic}_{lang}_{ind}",
            f"{base_topic}_{ind}_{lang}",
            f"{base_topic}_{lang}",
            f"{base_topic}_{ind}",
            base_topic
        ]

        best_file = None
        for root, _, files in os.walk(library_path):
            for cand in candidates:
                for f in files:
                    f_path = Path(root) / f
                    if f_path.stem == cand:
                        best_file = f_path
                        break
                if best_file:
                    break
            if best_file:
                break

        if best_file:
            found_files.append(best_file)
            solution_name = f"{best_file.stem}_solution{best_file.suffix}"
            solution_path = best_file.parent / solution_name
            if solution_path.exists():
                found_files.append(solution_path)

        return found_files

    @staticmethod
    def _find_screenshot(topic_folder: Path, shape_name: str, lang: str, ind: str) -> Optional[Path]:
        """
        Looks for a matching image in the 'screenshots' subfolder.
        """
        screenshots_dir = topic_folder / "screenshots"
        if not screenshots_dir.exists():
            return None

        extensions = ['.png', '.jpg', '.jpeg']
        candidates = [
            f"{shape_name}_{lang}_{ind}",
            f"{shape_name}_{ind}_{lang}",
            f"{shape_name}_{lang}",
            f"{shape_name}_{ind}",
            shape_name
        ]

        try:
            files = os.listdir(screenshots_dir)
            for cand in candidates:
                for ext in extensions:
                    target = f"{cand}{ext}"
                    for f in files:
                        if f.lower() == target.lower():
                            return screenshots_dir / f
        except OSError:
            pass
        return None

    @staticmethod
    def _replace_image_contain(slide, shape, image_path: Path) -> bool:
        """
        Replaces a shape with an image, maintaining aspect ratio (contain) and centering.
        """
        try:
            old_left, old_top = shape.left, shape.top
            old_width, old_height = shape.width, shape.height

            # Insert new image - python-pptx expects string or stream
            new_pic = slide.shapes.add_picture(str(image_path), 0, 0)

            native_width = new_pic.width
            native_height = new_pic.height
            if native_width == 0 or native_height == 0:
                return False

            # Calculate Scale
            aspect_img = native_width / native_height
            aspect_box = old_width / old_height

            if aspect_img > aspect_box:
                # Image is wider than box
                new_width = old_width
                new_height = int(old_width / aspect_img)
            else:
                # Image is taller than box
                new_height = old_height
                new_width = int(old_height * aspect_img)

            # Center
            offset_x = (old_width - new_width) // 2
            offset_y = (old_height - new_height) // 2

            new_pic.left = old_left + offset_x
            new_pic.top = old_top + offset_y
            new_pic.width = new_width
            new_pic.height = new_height

            # Remove placeholder
            sp = shape.element
            sp.getparent().remove(sp)
            return True
        except (AttributeError, ValueError, OSError) as e:
            print(f"ERROR replacing image: {e}")
            return False

    @staticmethod
    def _apply_text_replacements(text_frame, context: Dict[str, str]):
        """
        Replaces [% key %] placeholders in text.
        """
        if not text_frame or not text_frame.text:
            return
        text = text_frame.text

        # Fixed: Removed redundant escape for closing bracket
        pattern = re.compile(r'\[%\s*([\w\-]+)\s*%]')
        matches = pattern.findall(text)
        if not matches:
            return

        new_text = text
        replaced = False
        for key in matches:
            val = None
            if key in context:
                val = context[key]
            elif key.lower() in context:
                val = context[key.lower()]

            if val is not None:
                key_pattern = re.compile(r'\[%\s*' + re.escape(key) + r'\s*%]', re.IGNORECASE)
                new_text = key_pattern.sub(str(val), new_text)
                replaced = True

        if replaced:
            try:
                text_frame.text = new_text
            except (AttributeError, ValueError):
                pass

    @classmethod
    def generate_session(cls, req, library_path: Path, output_path: Path) -> Dict[str, Any]:
        """
        Main orchestration function.
        """
        # 1. Setup Folders
        folder_name = f"{req.date}_{req.customer_name}_{req.session_name}".replace(" ", "_")
        target_dir = output_path / folder_name
        target_dir.mkdir(parents=True, exist_ok=True)
        exercises_dir = target_dir / "exercises"
        exercises_dir.mkdir(exist_ok=True)

        print(f"INFO: Generating session for {req.customer_name}")

        # 2. System Variables
        system_vars = {
            "customer_name": req.customer_name,
            "customer": req.customer_name,
            "klant": req.customer_name,
            "session_name": req.session_name,
            "session": req.session_name,
            "date": req.date,
            "datum": req.date,
            "industry": req.customer_industry,
            "industry_code": req.industry_code,
            "language": req.language_code
        }

        files_copied = 0
        pptx_merge_list = []

        # 3. Gather Files
        # A. Intro
        intro_matches = cls.find_best_matches(library_path, "Intro", req.language_code, req.industry_code)
        if intro_matches:
            pptx_merge_list.append(intro_matches[0])

        # B. Sections & Topics
        for section in req.sections:
            for topic in section.topics:
                matches = cls.find_best_matches(library_path, topic, req.language_code, req.industry_code)
                for file_path in matches:
                    if file_path.suffix.lower() == '.pptx':
                        pptx_merge_list.append(file_path)
                    else:
                        # Copy Exercises with Section Prefix
                        dest_folder = exercises_dir
                        final_name = file_path.name
                        dest_path = dest_folder / final_name

                        if dest_path.exists():
                            safe_title = section.title.replace("/", "-").replace("\\", "-")
                            final_name = f"{safe_title} - {file_path.name}"
                            dest_path = dest_folder / final_name
                            counter = 1
                            while dest_path.exists():
                                final_name = f"{safe_title} - {file_path.stem}_{counter}{file_path.suffix}"
                                dest_path = dest_folder / final_name
                                counter += 1

                        try:
                            shutil.copy2(file_path, dest_path)
                            files_copied += 1
                        except OSError:
                            pass

        # C. Outro
        outro_matches = cls.find_best_matches(library_path, "Outro", req.language_code, req.industry_code)
        if outro_matches:
            pptx_merge_list.append(outro_matches[0])

        # 4. Merge & Process PPTX
        if pptx_merge_list:
            print(f"INFO: Merging {len(pptx_merge_list)} presentations...")

            master_path = pptx_merge_list[0]
            master_context = TranslationService.build_file_context(
                library_path, master_path, req.language_code, req.industry_code, system_vars
            )

            # python-pptx expects string path or file-like object
            prs = Presentation(str(master_path))

            # Process Master Slides
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        cls._apply_text_replacements(shape.text_frame, master_context)

                    img_path = cls._find_screenshot(master_path.parent, shape.name, req.language_code, req.industry_code)
                    if img_path:
                        cls._replace_image_contain(slide, shape, img_path)

            # Process Sub-Presentations
            for i in range(1, len(pptx_merge_list)):
                sub_path = pptx_merge_list[i]
                file_context = TranslationService.build_file_context(
                    library_path, sub_path, req.language_code, req.industry_code, system_vars
                )

                try:
                    sub_prs = Presentation(str(sub_path))
                    for slide in sub_prs.slides:
                        # Layout
                        layout_idx = sub_prs.slide_layouts.index(slide.slide_layout)
                        try:
                            slide_layout = prs.slide_layouts[layout_idx]
                        except IndexError:
                            slide_layout = prs.slide_layouts[0]

                        new_slide = prs.slides.add_slide(slide_layout)

                        for shape in slide.shapes:
                            # Placeholders
                            if shape.is_placeholder:
                                try:
                                    ph_idx = shape.placeholder_format.idx
                                    new_ph = new_slide.placeholders[ph_idx]
                                    if shape.has_text_frame:
                                        new_ph.text = shape.text_frame.text
                                        cls._apply_text_replacements(new_ph.text_frame, file_context)
                                except KeyError:
                                    pass

                            # Normal Shapes
                            else:
                                img_match = cls._find_screenshot(sub_path.parent, shape.name, req.language_code, req.industry_code)
                                if img_match:
                                    # Create temp picture to swap
                                    temp_pic = new_slide.shapes.add_picture(str(img_match), shape.left, shape.top, shape.width, shape.height)
                                    cls._replace_image_contain(new_slide, temp_pic, img_match)

                                elif shape.has_text_frame:
                                    new_shape = new_slide.shapes.add_textbox(
                                        shape.left, shape.top, shape.width, shape.height
                                    )
                                    new_shape.text = shape.text
                                    cls._apply_text_replacements(new_shape.text_frame, file_context)

                except Exception as e:
                    # Narrowed for specific file processing errors
                    print(f"ERROR processing {sub_path.name}: {e}")

            # Save - python-pptx save() accepts string path
            output_pptx = target_dir / "slides.pptx"
            prs.save(str(output_pptx))
            files_copied += 1

        return {
            "status": "success",
            "target_dir": str(target_dir),
            "files_count": files_copied
        }